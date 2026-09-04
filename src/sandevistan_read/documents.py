from __future__ import annotations

import hashlib
import html
import base64
import io
import os
import posixpath
import re
import shutil
import subprocess
import uuid
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from urllib.parse import unquote
from typing import Any
from xml.etree import ElementTree

import fitz
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from PIL import Image

from .config import CONFIG
from .paths import PATHS


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".epub", ".txt", ".md", ".markdown", ".html", ".htm", ".png", ".jpg", ".jpeg", ".webp"}
EPUB_MAX_ENTRIES = 10_000
EPUB_MAX_UNCOMPRESSED_BYTES = 1024 * 1024 * 1024
EPUB_MAX_TEXT_MEMBER_BYTES = 32 * 1024 * 1024


@dataclass
class ParsedBlock:
    text: str
    locator: dict[str, Any]
    image_path: str | None = None
    visual_needed: bool = False


@dataclass
class ParsedDocument:
    blocks: list[ParsedBlock] = field(default_factory=list)
    page_count: int = 0
    parser: str = ""
    preview_path: str | None = None
    metadata: dict[str, Any] = field(default_factory=dict)


def sanitize_filename(filename: str) -> str:
    filename = Path(filename).name.replace("\x00", "")
    cleaned = re.sub(r"[^\w.()\-\u4e00-\u9fff ]+", "_", filename, flags=re.UNICODE).strip(" .")
    return cleaned[:180] or "document"


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _clean_text(value: str) -> str:
    value = value.replace("\x00", " ")
    value = re.sub(r"[ \t]+", " ", value)
    value = re.sub(r"\n{3,}", "\n\n", value)
    return value.strip()


def _render_pdf_page(page: fitz.Page, destination: Path) -> str:
    destination.parent.mkdir(parents=True, exist_ok=True)
    pixmap = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
    pixmap.save(destination)
    return str(destination.relative_to(PATHS.root))


def _store_visual(data: bytes, destination: Path) -> str:
    destination.parent.mkdir(parents=True, exist_ok=True)
    try:
        with Image.open(io.BytesIO(data)) as image:
            image.convert("RGB").save(destination, "PNG")
    except Exception as exc:
        raise ValueError("图片格式无法解码") from exc
    return str(destination.relative_to(PATHS.root))


def _store_svg(data: bytes, destination: Path) -> str:
    destination.parent.mkdir(parents=True, exist_ok=True)
    try:
        if len(data) > 16 * 1024 * 1024 or re.search(br"<!DOCTYPE|<script|<foreignObject", data, re.I):
            raise ValueError("SVG 包含不允许的内容")
        if re.search(br"(?:href|xlink:href)\s*=\s*['\"](?!data:|#)", data, re.I):
            raise ValueError("SVG 不允许引用外部资源")
        with fitz.open(stream=data, filetype="svg") as document:
            pixmap = document[0].get_pixmap(matrix=fitz.Matrix(1.5, 1.5), alpha=False)
            pixmap.save(destination)
    except Exception as exc:
        raise ValueError("SVG 无法安全渲染") from exc
    return str(destination.relative_to(PATHS.root))


def parse_pdf(path: Path, source_id: str) -> ParsedDocument:
    document = fitz.open(path)
    result = ParsedDocument(page_count=len(document), parser="pymupdf", preview_path=str(path.relative_to(PATHS.root)), metadata={"locator_unit": "page"})
    render_dir = PATHS.renders / source_id
    for index, page in enumerate(document):
        blocks = page.get_text("blocks", sort=True)
        text = _clean_text("\n".join(str(block[4]) for block in blocks if len(block) > 4))
        image_count = len(page.get_images(full=True))
        drawing_count = len(page.get_drawings())
        page_area = max(page.rect.width * page.rect.height, 1)
        text_density = len(text) / page_area
        visual_needed = len(text) < 80 or image_count > 0 or drawing_count > 0 or text_density < 0.00018
        image_path = None
        if visual_needed:
            image_path = _render_pdf_page(page, render_dir / f"page-{index + 1:04d}.png")
        if text:
            locator = {
                "kind": "page",
                "page": index + 1,
                "bboxes": [list(map(float, block[:4])) for block in blocks[:24]],
            }
            result.blocks.append(ParsedBlock(text=text, locator=locator, image_path=image_path, visual_needed=visual_needed))
        elif image_path:
            result.blocks.append(
                ParsedBlock(text="", locator={"kind": "page", "page": index + 1, "bboxes": []}, image_path=image_path, visual_needed=True)
            )
    document.close()
    return result


def _convert_office_to_pdf(path: Path, source_id: str) -> Path | None:
    executable = CONFIG.tools.libreoffice_path
    if not executable:
        return None
    output_dir = PATHS.renders / source_id / "office-preview"
    output_dir.mkdir(parents=True, exist_ok=True)
    profile = PATHS.libreoffice_profiles / f"{source_id}-{uuid.uuid4().hex}"
    profile.mkdir(parents=True, exist_ok=False)
    command = [
        executable,
        f"-env:UserInstallation={profile.as_uri()}",
        "--headless",
        "--nologo",
        "--nodefault",
        "--nofirststartwizard",
        "--nolockcheck",
        "--convert-to",
        "pdf",
        "--outdir",
        str(output_dir),
        str(path),
    ]
    environment = os.environ.copy()
    environment["SAL_USE_VCLPLUGIN"] = "svp"
    try:
        completed = subprocess.run(command, capture_output=True, text=True, timeout=180, check=False, env=environment)
        converted = output_dir / f"{path.stem}.pdf"
        return converted if completed.returncode == 0 and converted.exists() else None
    finally:
        shutil.rmtree(profile, ignore_errors=True)


def parse_docx(path: Path, source_id: str) -> ParsedDocument:
    document = Document(path)
    blocks: list[ParsedBlock] = []
    section = "文档"
    ordinal = 0
    for paragraph in document.paragraphs:
        text = _clean_text(paragraph.text)
        if not text:
            continue
        style = paragraph.style.name.lower() if paragraph.style else ""
        if "heading" in style or "标题" in style:
            section = text
        blocks.append(ParsedBlock(text=text, locator={"kind": "section", "section": section, "paragraph": ordinal + 1}))
        ordinal += 1
    for table_index, table in enumerate(document.tables, start=1):
        rows = [" | ".join(_clean_text(cell.text) for cell in row.cells) for row in table.rows]
        table_text = _clean_text("\n".join(rows))
        if table_text:
            blocks.append(ParsedBlock(text=table_text, locator={"kind": "table", "table": table_index, "section": section}))
    preview = _convert_office_to_pdf(path, source_id)
    page_count = 0
    if preview:
        with fitz.open(preview) as pdf:
            page_count = len(pdf)
            render_dir = PATHS.renders / source_id
            for index, page in enumerate(pdf):
                image_path = _render_pdf_page(page, render_dir / f"page-{index + 1:04d}.png")
                blocks.append(ParsedBlock(text="", locator={"kind": "page", "page": index + 1, "visual_only": True}, image_path=image_path, visual_needed=True))
    return ParsedDocument(
        blocks=blocks,
        page_count=page_count or max(1, len(document.sections)),
        parser="python-docx+libreoffice" if preview else "python-docx",
        preview_path=str(preview.relative_to(PATHS.root)) if preview else None,
        metadata={"visual_preview": bool(preview), "locator_unit": "section"},
    )


def parse_pptx(path: Path, source_id: str) -> ParsedDocument:
    presentation = Presentation(path)
    blocks: list[ParsedBlock] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        image_count = 0
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                value = _clean_text(shape.text)
                if value:
                    parts.append(value)
            if getattr(shape, "shape_type", None) == 13:
                image_count += 1
            if getattr(shape, "has_table", False):
                rows = [" | ".join(_clean_text(cell.text) for cell in row.cells) for row in shape.table.rows]
                table_text = _clean_text("\n".join(rows))
                if table_text:
                    parts.append("[表格]\n" + table_text)
        text = _clean_text("\n".join(parts))
        blocks.append(
            ParsedBlock(
                text=text,
                locator={"kind": "slide", "slide": slide_number},
                visual_needed=image_count > 0 or len(text) < 80,
            )
        )
    preview = _convert_office_to_pdf(path, source_id)
    if preview:
        with fitz.open(preview) as pdf:
            render_dir = PATHS.renders / source_id
            for index, page in enumerate(pdf):
                if index < len(blocks) and blocks[index].visual_needed:
                    blocks[index].image_path = _render_pdf_page(page, render_dir / f"slide-{index + 1:04d}.png")
    return ParsedDocument(
        blocks=blocks,
        page_count=len(presentation.slides),
        parser="python-pptx+libreoffice" if preview else "python-pptx",
        preview_path=str(preview.relative_to(PATHS.root)) if preview else None,
        metadata={"visual_preview": bool(preview), "locator_unit": "slide"},
    )


def _epub_member(base: str, href: str) -> str:
    href = unquote(href.split("#", 1)[0]).replace("\\", "/")
    member = posixpath.normpath(posixpath.join(base, href))
    if not member or member.startswith("../") or member.startswith("/") or "/../" in f"/{member}/":
        raise ValueError("EPUB contains an unsafe resource path")
    return member


def _xml_local_name(tag: str) -> str:
    return tag.rsplit("}", 1)[-1]


def parse_epub(path: Path, source_id: str = "preview") -> ParsedDocument:
    try:
        archive = zipfile.ZipFile(path)
    except zipfile.BadZipFile as exc:
        raise ValueError("EPUB 文件损坏或不是有效 ZIP 容器") from exc
    with archive:
        entries = archive.infolist()
        if len(entries) > EPUB_MAX_ENTRIES or sum(item.file_size for item in entries) > EPUB_MAX_UNCOMPRESSED_BYTES:
            raise ValueError("EPUB 解压后体积或文件数量超过安全限制")
        names = {item.filename for item in entries}
        if "META-INF/container.xml" not in names:
            raise ValueError("EPUB 缺少 META-INF/container.xml")
        try:
            container = ElementTree.fromstring(archive.read("META-INF/container.xml"))
            rootfiles = [node.attrib.get("full-path", "") for node in container.iter() if _xml_local_name(node.tag) == "rootfile"]
        except ElementTree.ParseError as exc:
            raise ValueError("EPUB container.xml 无法解析") from exc
        opf_name = next((name for name in rootfiles if name in names), "")
        if not opf_name:
            raise ValueError("EPUB 未找到 package document")
        try:
            package = ElementTree.fromstring(archive.read(opf_name))
        except ElementTree.ParseError as exc:
            raise ValueError("EPUB package document 无法解析") from exc
        base = posixpath.dirname(opf_name)
        manifest: dict[str, dict[str, str]] = {}
        spine: list[str] = []
        metadata: dict[str, Any] = {"locator_unit": "chapter"}
        creators: list[str] = []
        for node in package.iter():
            local = _xml_local_name(node.tag)
            if local == "item" and node.attrib.get("id"):
                manifest[node.attrib["id"]] = dict(node.attrib)
            elif local == "itemref" and node.attrib.get("idref"):
                spine.append(node.attrib["idref"])
            elif local in {"title", "language"} and (node.text or "").strip() and local not in metadata:
                metadata[local] = _clean_text(node.text or "")
            elif local == "creator" and (node.text or "").strip():
                creators.append(_clean_text(node.text or ""))
        if creators:
            metadata["creators"] = creators
        encrypted: set[str] = set()
        if "META-INF/encryption.xml" in names:
            try:
                encryption = ElementTree.fromstring(archive.read("META-INF/encryption.xml"))
                encrypted = {unquote(node.attrib.get("URI", "")) for node in encryption.iter() if _xml_local_name(node.tag) == "CipherReference"}
            except ElementTree.ParseError:
                encrypted = set()
        blocks: list[ParsedBlock] = []
        for spine_index, item_id in enumerate(spine, start=1):
            item = manifest.get(item_id, {})
            href = item.get("href", "")
            media_type = item.get("media-type", "")
            if not href or media_type not in {"application/xhtml+xml", "text/html"}:
                continue
            member = _epub_member(base, href)
            if member not in names:
                continue
            if href in encrypted or member in encrypted:
                raise ValueError("EPUB 正文受 DRM/加密保护，无法在本地解析")
            info = archive.getinfo(member)
            if info.file_size > EPUB_MAX_TEXT_MEMBER_BYTES:
                raise ValueError("EPUB 单个正文文件超过安全限制")
            raw = archive.read(member).decode("utf-8", errors="replace")
            soup = BeautifulSoup(raw, "html.parser")
            for image_index, image in enumerate(soup.find_all("img"), start=1):
                src = str(image.get("src") or "")
                if not src or src.startswith(("http://", "https://")):
                    continue
                try:
                    image_member = _epub_member(posixpath.dirname(member), src)
                    if image_member not in names or image_member in encrypted:
                        continue
                    image_data = archive.read(image_member)
                    destination = PATHS.renders / source_id / f"epub-{spine_index:04d}-{image_index:04d}.png"
                    stored = _store_svg(image_data, destination) if image_member.lower().endswith(".svg") else _store_visual(image_data, destination)
                    blocks.append(ParsedBlock(
                        text="", image_path=stored, visual_needed=True,
                        locator={"kind": "epub-image", "spine": spine_index, "href": href, "asset": src, "visual_only": True},
                    ))
                except (KeyError, ValueError):
                    continue
            for svg_index, svg in enumerate(soup.find_all("svg"), start=1):
                try:
                    stored = _store_svg(str(svg).encode(), PATHS.renders / source_id / f"epub-svg-{spine_index:04d}-{svg_index:04d}.png")
                    blocks.append(ParsedBlock("", {"kind": "epub-svg", "spine": spine_index, "href": href, "visual_only": True}, stored, True))
                except ValueError:
                    continue
            for unsafe in soup(["script", "style", "iframe", "object", "embed", "svg"]):
                unsafe.decompose()
            heading = soup.find(re.compile(r"^h[1-6]$"))
            title = soup.title.get_text(" ", strip=True) if soup.title else ""
            section = _clean_text(heading.get_text(" ", strip=True) if heading else title) or f"章节 {spine_index}"
            text = _clean_text(soup.get_text("\n"))
            if text:
                blocks.append(ParsedBlock(text=text, locator={"kind": "epub", "spine": spine_index, "href": href, "section": section}))
        if not blocks:
            raise ValueError("EPUB 没有可读取的正文")
        metadata["spine_items"] = len(spine)
        try:
            preview_path = str(path.relative_to(PATHS.root))
        except ValueError:
            preview_path = None
        return ParsedDocument(blocks=blocks, page_count=len(blocks), parser="epub-spine", preview_path=preview_path, metadata=metadata)


def parse_text(path: Path, extension: str, source_id: str = "preview") -> ParsedDocument:
    raw = path.read_text(encoding="utf-8", errors="replace")
    visual_blocks: list[ParsedBlock] = []
    if extension in {".html", ".htm", ".md", ".markdown"}:
        pattern = re.compile(r"data:(image/(?:png|jpeg|webp));base64,([A-Za-z0-9+/=\s]+)", re.I)
        for index, match in enumerate(pattern.finditer(raw), start=1):
            try:
                data = base64.b64decode(re.sub(r"\s+", "", match.group(2)), validate=True)
                stored = _store_visual(data, PATHS.renders / source_id / f"inline-{index:04d}.png")
                line = raw.count("\n", 0, match.start()) + 1
                visual_blocks.append(ParsedBlock("", {"kind": "inline-image", "line": line, "visual_only": True}, stored, True))
            except (ValueError, base64.binascii.Error):
                continue
        if extension in {".html", ".htm"}:
            inline_soup = BeautifulSoup(raw, "html.parser")
            for index, svg in enumerate(inline_soup.find_all("svg"), start=1):
                try:
                    stored = _store_svg(str(svg).encode(), PATHS.renders / source_id / f"inline-svg-{index:04d}.png")
                    visual_blocks.append(ParsedBlock("", {"kind": "inline-svg", "visual_only": True}, stored, True))
                except ValueError:
                    continue
    if extension in {".html", ".htm"}:
        soup = BeautifulSoup(raw, "html.parser")
        for unsafe in soup(["script", "style", "iframe", "object", "embed"]):
            unsafe.decompose()
        raw = soup.get_text("\n")
    lines = raw.splitlines()
    blocks: list[ParsedBlock] = []
    buffer: list[str] = []
    start_line = 1
    section = "文档"
    for line_number, line in enumerate(lines, start=1):
        stripped = line.strip()
        if extension in {".md", ".markdown"} and stripped.startswith("#"):
            if buffer:
                blocks.append(ParsedBlock(_clean_text("\n".join(buffer)), {"kind": "lines", "line_start": start_line, "line_end": line_number - 1, "section": section}))
                buffer = []
            section = stripped.lstrip("# ") or section
            start_line = line_number
        if not buffer:
            start_line = line_number
        buffer.append(line)
        if sum(len(item) for item in buffer) >= 2200:
            blocks.append(ParsedBlock(_clean_text("\n".join(buffer)), {"kind": "lines", "line_start": start_line, "line_end": line_number, "section": section}))
            buffer = []
    if buffer:
        blocks.append(ParsedBlock(_clean_text("\n".join(buffer)), {"kind": "lines", "line_start": start_line, "line_end": len(lines), "section": section}))
    return ParsedDocument(blocks=[block for block in blocks if block.text] + visual_blocks, page_count=max(1, len(blocks)), parser=f"text-{extension.lstrip('.')}", metadata={"locator_unit": "section", "inline_visuals": len(visual_blocks)})


def parse_image(path: Path, source_id: str) -> ParsedDocument:
    stored = _store_visual(path.read_bytes(), PATHS.renders / source_id / "image-0001.png")
    return ParsedDocument(
        blocks=[ParsedBlock("", {"kind": "image", "visual_only": True}, stored, True)],
        page_count=1,
        parser="pymupdf-image",
        preview_path=stored,
        metadata={"locator_unit": "image"},
    )


def parse_document(path: Path, source_id: str) -> ParsedDocument:
    extension = path.suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Unsupported document type: {extension}")
    if extension == ".pdf":
        return parse_pdf(path, source_id)
    if extension == ".docx":
        return parse_docx(path, source_id)
    if extension == ".pptx":
        return parse_pptx(path, source_id)
    if extension == ".epub":
        return parse_epub(path, source_id)
    if extension in {".png", ".jpg", ".jpeg", ".webp"}:
        return parse_image(path, source_id)
    return parse_text(path, extension, source_id)


def chunk_blocks(blocks: list[ParsedBlock], target_chars: int = 1800, overlap_chars: int = 260) -> list[ParsedBlock]:
    chunks: list[ParsedBlock] = []
    for block in blocks:
        text = _clean_text(block.text)
        if not text:
            continue
        if len(text) <= target_chars:
            chunks.append(block)
            continue
        start = 0
        while start < len(text):
            end = min(len(text), start + target_chars)
            if end < len(text):
                boundary = max(text.rfind("。", start, end), text.rfind("\n", start, end), text.rfind(". ", start, end))
                if boundary > start + target_chars // 2:
                    end = boundary + 1
            chunks.append(ParsedBlock(text=text[start:end].strip(), locator=dict(block.locator), image_path=block.image_path, visual_needed=False))
            if end >= len(text):
                break
            start = max(end - overlap_chars, start + 1)
    return chunks
